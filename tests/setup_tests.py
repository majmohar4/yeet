#!/usr/bin/env python3
"""Generate all test assets needed by the test suite."""

import os
import zipfile
from pathlib import Path


def setup_test_environment():
    test_dir = Path('test_assets')
    test_dir.mkdir(exist_ok=True)
    print("Generating test assets...")

    # ── JPEG ──────────────────────────────────────────────────────────────────
    try:
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (800, 600), color=(73, 109, 137))
        draw = ImageDraw.Draw(img)
        draw.text((400, 300), "Test Image", fill=(255, 255, 255), anchor="mm")
        img.save(test_dir / 'test_image.jpg', 'JPEG')
        Image.new('RGB', (400, 300), color=(200, 100, 50)).save(test_dir / 'test_image.png', 'PNG')
        print("  ✓ test_image.jpg, test_image.png")
    except ImportError:
        # Minimal 1×1 JPEG/PNG without Pillow
        _write_minimal_jpeg(test_dir / 'test_image.jpg')
        _write_minimal_png(test_dir / 'test_image.png')
        print("  ✓ minimal test_image.jpg, test_image.png (no Pillow)")

    # ── PDF ───────────────────────────────────────────────────────────────────
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        c = canvas.Canvas(str(test_dir / 'test.pdf'), pagesize=letter)
        c.drawString(100, 750, "Test PDF Document")
        c.save()
        print("  ✓ test.pdf (reportlab)")
    except ImportError:
        (test_dir / 'test.pdf').write_bytes(_minimal_pdf())
        print("  ✓ test.pdf (minimal)")

    # ── DOCX ──────────────────────────────────────────────────────────────────
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('Test DOCX for Yeet.')
        doc.save(test_dir / 'test.docx')
        print("  ✓ test.docx (python-docx)")
    except ImportError:
        _write_minimal_office_zip(test_dir / 'test.docx')
        print("  ✓ test.docx (minimal zip)")

    # ── XLSX ──────────────────────────────────────────────────────────────────
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws['A1'] = 'Test'
        ws['B1'] = 'Data'
        wb.save(test_dir / 'test.xlsx')
        print("  ✓ test.xlsx (openpyxl)")
    except ImportError:
        _write_minimal_office_zip(test_dir / 'test.xlsx')
        print("  ✓ test.xlsx (minimal zip)")

    # ── PPTX ──────────────────────────────────────────────────────────────────
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"
        prs.save(test_dir / 'test.pptx')
        print("  ✓ test.pptx (python-pptx)")
    except ImportError:
        _write_minimal_office_zip(test_dir / 'test.pptx')
        print("  ✓ test.pptx (minimal zip)")

    # ── ZIP ───────────────────────────────────────────────────────────────────
    with zipfile.ZipFile(test_dir / 'test.zip', 'w') as zf:
        zf.writestr('readme.txt', 'This is a test ZIP file')
        zf.writestr('data.txt', 'Sample data')
    print("  ✓ test.zip")

    # ── Large file (50 MB) ────────────────────────────────────────────────────
    large = test_dir / 'large_50mb.bin'
    with open(large, 'wb') as f:
        chunk = os.urandom(1024 * 1024)
        for _ in range(50):
            f.write(chunk)
    print("  ✓ large_50mb.bin")

    # ── EICAR test virus ──────────────────────────────────────────────────────
    eicar = b'X5O!P%@AP[4\\PZX54(P^)7CC)7}$EICAR-STANDARD-ANTIVIRUS-TEST-FILE!$H+H*'
    (test_dir / 'eicar.txt').write_bytes(eicar)
    print("  ✓ eicar.txt (ClamAV test file)")

    # ── Text variants ─────────────────────────────────────────────────────────
    (test_dir / 'empty.txt').write_text('')
    (test_dir / 'file_with_spaces.txt').write_text('Spaces in filename test')
    (test_dir / 'unicode_test.txt').write_text('Unicode content: こんにちは', encoding='utf-8')
    print("  ✓ text file variants")

    # ── Blocked files (should be rejected by server) ──────────────────────────
    blocked_dir = test_dir / 'blocked'
    blocked_dir.mkdir(exist_ok=True)
    (blocked_dir / 'malware.exe').write_bytes(b'MZ\x90\x00test')
    (blocked_dir / 'script.sh').write_text('#!/bin/bash\necho pwned')
    (blocked_dir / 'xss.html').write_text('<script>alert(1)</script>')
    (blocked_dir / 'payload.js').write_text('alert(1)')
    print("  ✓ blocked/ files (for rejection tests)")

    print(f"\n✓ All test assets in: {test_dir.resolve()}\n")


def _minimal_pdf() -> bytes:
    return b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj
3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj
xref
0 4
0000000000 65535 f\x20
0000000009 00000 n\x20
0000000058 00000 n\x20
0000000115 00000 n\x20
trailer<</Size 4/Root 1 0 R>>
startxref
203
%%EOF"""


def _write_minimal_office_zip(path: Path) -> None:
    with zipfile.ZipFile(path, 'w') as zf:
        zf.writestr('[Content_Types].xml',
                    '<?xml version="1.0"?>'
                    '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"/>')
        zf.writestr('_rels/.rels',
                    '<?xml version="1.0"?>'
                    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>')


def _write_minimal_jpeg(path: Path) -> None:
    # 1×1 white JPEG
    data = bytes([
        0xFF, 0xD8, 0xFF, 0xE0, 0x00, 0x10, 0x4A, 0x46, 0x49, 0x46, 0x00, 0x01,
        0x01, 0x00, 0x00, 0x01, 0x00, 0x01, 0x00, 0x00, 0xFF, 0xDB, 0x00, 0x43,
        0x00, 0x08, 0x06, 0x06, 0x07, 0x06, 0x05, 0x08, 0x07, 0x07, 0x07, 0x09,
        0x09, 0x08, 0x0A, 0x0C, 0x14, 0x0D, 0x0C, 0x0B, 0x0B, 0x0C, 0x19, 0x12,
        0x13, 0x0F, 0x14, 0x1D, 0x1A, 0x1F, 0x1E, 0x1D, 0x1A, 0x1C, 0x1C, 0x20,
        0x24, 0x2E, 0x27, 0x20, 0x22, 0x2C, 0x23, 0x1C, 0x1C, 0x28, 0x37, 0x29,
        0x2C, 0x30, 0x31, 0x34, 0x34, 0x34, 0x1F, 0x27, 0x39, 0x3D, 0x38, 0x32,
        0x3C, 0x2E, 0x33, 0x34, 0x32, 0xFF, 0xC0, 0x00, 0x0B, 0x08, 0x00, 0x01,
        0x00, 0x01, 0x01, 0x01, 0x11, 0x00, 0xFF, 0xC4, 0x00, 0x1F, 0x00, 0x00,
        0x01, 0x05, 0x01, 0x01, 0x01, 0x01, 0x01, 0x01, 0x00, 0x00, 0x00, 0x00,
        0x00, 0x00, 0x00, 0x00, 0x01, 0x02, 0x03, 0x04, 0x05, 0x06, 0x07, 0x08,
        0x09, 0x0A, 0x0B, 0xFF, 0xC4, 0x00, 0xB5, 0x10, 0x00, 0x02, 0x01, 0x03,
        0x03, 0x02, 0x04, 0x03, 0x05, 0x05, 0x04, 0x04, 0x00, 0x00, 0x01, 0x7D,
        0x01, 0x02, 0x03, 0x00, 0x04, 0x11, 0x05, 0x12, 0x21, 0x31, 0x41, 0x06,
        0x13, 0x51, 0x61, 0x07, 0x22, 0x71, 0x14, 0x32, 0x81, 0x91, 0xA1, 0x08,
        0x23, 0x42, 0xB1, 0xC1, 0x15, 0x52, 0xD1, 0xF0, 0x24, 0x33, 0x62, 0x72,
        0x82, 0x09, 0x0A, 0x16, 0x17, 0x18, 0x19, 0x1A, 0x25, 0x26, 0x27, 0x28,
        0x29, 0x2A, 0x34, 0x35, 0x36, 0x37, 0x38, 0x39, 0x3A, 0x43, 0x44, 0x45,
        0x46, 0x47, 0x48, 0x49, 0x4A, 0x53, 0x54, 0x55, 0x56, 0x57, 0x58, 0x59,
        0x5A, 0x63, 0x64, 0x65, 0x66, 0x67, 0x68, 0x69, 0x6A, 0x73, 0x74, 0x75,
        0x76, 0x77, 0x78, 0x79, 0x7A, 0x83, 0x84, 0x85, 0x86, 0x87, 0x88, 0x89,
        0x8A, 0x92, 0x93, 0x94, 0x95, 0x96, 0x97, 0x98, 0x99, 0x9A, 0xA2, 0xA3,
        0xA4, 0xA5, 0xA6, 0xA7, 0xA8, 0xA9, 0xAA, 0xB2, 0xB3, 0xB4, 0xB5, 0xB6,
        0xB7, 0xB8, 0xB9, 0xBA, 0xC2, 0xC3, 0xC4, 0xC5, 0xC6, 0xC7, 0xC8, 0xC9,
        0xCA, 0xD2, 0xD3, 0xD4, 0xD5, 0xD6, 0xD7, 0xD8, 0xD9, 0xDA, 0xE1, 0xE2,
        0xE3, 0xE4, 0xE5, 0xE6, 0xE7, 0xE8, 0xE9, 0xEA, 0xF1, 0xF2, 0xF3, 0xF4,
        0xF5, 0xF6, 0xF7, 0xF8, 0xF9, 0xFA, 0xFF, 0xDA, 0x00, 0x08, 0x01, 0x01,
        0x00, 0x00, 0x3F, 0x00, 0xFB, 0xD4, 0xFF, 0xD9,
    ])
    path.write_bytes(data)


def _write_minimal_png(path: Path) -> None:
    # 1×1 red PNG
    import struct, zlib
    def chunk(name, data):
        c = name + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xFFFFFFFF)
    sig = b'\x89PNG\r\n\x1a\n'
    ihdr = chunk(b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0))
    raw = b'\x00\xFF\x00\x00'  # filter byte + RGB
    idat = chunk(b'IDAT', zlib.compress(raw))
    iend = chunk(b'IEND', b'')
    path.write_bytes(sig + ihdr + idat + iend)


if __name__ == '__main__':
    setup_test_environment()
